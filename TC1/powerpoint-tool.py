import pandas as pd
from pptx import Presentation

PATH = "./4Q/4Q24 Webinar Statistics.xlsx"
TEMPLATE_PATH = "./CatRock_Template.pptx"
OUTPUT_PATH = "./4Q/Cat Rock Capital 4Q24 Review Webinar Presentation.pptx"

# LOADING DATA FOR SLIDE 4

t4 = pd.read_excel(PATH, sheet_name=0, header=2)
t4.columns = ["Stock", "3Q24", "YTD 3Q24", "Past 12 Months"] + list(t4.columns[4:])

t4 = t4[t4["Stock"].isin(["Cat Rock (Net)", "S&P 500", "MSCI World"])]
t4 = t4[["Stock", "3Q24", "YTD 3Q24", "Past 12 Months"]]

t4_clean = t4.iloc[:, [0, 1, 3, 4]].copy()
t4_clean.columns = ["Stock", "3Q24", "YTD 3Q24", "Past 12 Months"]

print("\nSlide 4: Cat Rock Results\n" + "-"*40)
print(t4_clean.to_string(index=False))

print("After 4...")

# LOADING DATA FOR SLIDE 5

t5 = pd.read_excel(PATH, sheet_name=1, header=4)
t5.columns = [c.strip() for c in t5.columns]
t5["Ticker"] = t5["Stock"].str.split().str[0]

t5 = t5.rename(columns={
    "% Change": "YTD Stock Price Performance",
    "% Contribution (Net)": "YTD Net Attribution"
})

wanted = ["META US EQUITY", "KSPI US EQUITY", "KSPI LI EQUITY", "GOOG US EQUITY", "ARES US EQUITY",
          "SEMR US EQUITY", "SCT LN EQUITY", "MSFT US EQUITY", "EVO SS EQUITY",
          "CTT AU EQUITY", "Other", "Total"]

t5 = t5[t5["Stock"].isin(wanted)]
t5 = t5[["Stock", "Ticker", "YTD Stock Price Performance", "YTD Net Attribution"]]

print("\nSlide 5: YTD Attribution\n" + "-"*40)
print(t5.to_string(index=False))

print("Loaded slide 5 data...")


# GENERATE PPTX

prs = Presentation(TEMPLATE_PATH)

slide4 = prs.slides[3]
shapes4 = slide4.shapes

for shape in shapes4:
    if (shape.has_table):
        table = shape.table
        print(len(table.rows), len(table.columns))
        for i in range(len(table.rows)):
            if i <= 1: continue
            for j in range(len(table.columns)):
                data = t4_clean.iloc[i-2,j]
                if isinstance(data, (int, float)): data = f"{abs(data*100):.1f}%"
                table.cell(i,j).text = str(data)

slide5 = prs.slides[4]
shapes5 = slide5.shapes

for shape in shapes5:
    if (shape.has_table):
        table = shape.table
        print(len(table.rows), len(table.columns))
        for i in range(len(table.rows)):
            if i <= 0: continue
            for j in range(len(table.columns)):
                data = t5.iloc[i-1,j]
                if isinstance(data, (int, float)): data = f"{abs(data*100):.1f}%"
                table.cell(i,j).text = str(data)


prs.save(OUTPUT_PATH)